from pptx import Presentation
from pptx.dml.color import RGBColor

from datetime import datetime
import json
import time
import glob
import os
from export_funcs import edit_title, plot_slide, transitn_slide, addtxtbox

#inputs accepted:
#1. template ppt file (needed as layouts are indexed from slide master)
#2. slide titles (transition, plots) (pass as list in order)
## list order: 
#[main,transition,probe plots...,transition,FD plots..,close]
#3. Image files list
##list order: [probe1,probe2,...,FD1,FD2..]
#3. # of probe plots, # of FD plots
#4. headers and comments for plot slides
#5. final ppt file save

def exptppt(jsonfile,outlier_json = None,stats_json = None,imagedir = None): 
    start = time.time()

    #Json Data Parsing
    f_1 = open(jsonfile,'r')
    data = json.loads(f_1.read())

    if outlier_json != None:
        f_2 = open(outlier_json,'r')
        outlier_data = json.loads(f_2.read())
    if stats_json != None:
        f_3 = open(stats_json,'r')
        stats_data = json.loads(f_3.read())

    #Query directory for png file names
    files = glob.glob(imagedir + "*.png")

    #Query + Cleanup lotid:WaferId data for outliers comments
    outlier_lst = [] #List of strings, string format = LotID: WaferId1, WaferId2,...
    for lotid,waferids in outlier_data.items():
        if len(waferids) > 1:
            one_line = lotid + ': '
            for waferid in waferids:
                if waferid == waferids[-1]:
                    one_line += waferid 
                    outlier_lst.append(one_line)
                else:
                    one_line += waferid + ','
        else:
            outlier_lst.append('{}: {}'.format(lotid,waferids[0]))
    
    outliers = '' #just 1 string = LotId: WaferId1, WaferId2,...
    for id in outlier_lst:
        if id == outlier_lst[-1]:
            outliers += id
        else:
            outliers += id + ', '
    
    #Create Presentation
    prs = Presentation(data[0]["image_dir"] + data[0]["template"])

    #Main Page, hardcoded(rarely changed)
    main_slide = prs.slides[0]
    date = datetime.now()
    edit_title(main_slide,data[0]["title"])
    subtitle = main_slide.shapes.placeholders[25]
    subtitle.text = date.strftime("%m/%d/%Y, %H:%M:%S")

    #Main page color defaults (hardcoded change when needed)
    for shape in main_slide.shapes:
        if shape.shape_id == 2:
            shape.fill.background()
        if shape.shape_id == 9: 
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(150,50,100)
        if shape.shape_id == 10: 
            shape.fill.background()
        if shape.shape_id == 11: 
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(100,100,100)


    x = 2
    #First transition and Probe plots
    transitn_slide(prs, data[1]["title"])

    for i in range(data[0]["probe"]):
        pngimage = [dir for dir in files if data[x]['title'] in dir]
        #if x == 2: #hardcode for this loop (outlier comments)
            #plot_slide(prs,header=data[x]["header"],comments=outliers,title=data[x]["title"],image=pngimage[0])
        #else:
        plot_slide(prs,header=data[x]["header"],comments=outliers,title=data[x]["title"],image=pngimage[0])
        x += 1

        #else:
        #    plot_slide(prs,data[x]["header"],data[x]["comment"],data[x]["title"],pngimage[0])
        #    x += 1
    #Second transition, FD plots, and last slide
    # transitn_slide(prs,data[x]["title"])

    # j = 1
    # for i in range(data[0]["FD"]):
    #     try:
    #         pngimage = [dir for dir in files if 'KDM' in dir]
    #         plot_slide(prs,header=data[x+i+j]["header"],comments=data[x+i+j]["comment"],title=data[x+i+j]["title"],image=pngimage[0])
    #     except:
    #         quit()

    layout = prs.slide_layouts[17]
    prs.slides.add_slide(layout)
    prs.save("json_final.pptx")
    
    end = time.time()
    print("Runtime: {} seconds".format(end-start))


#exptppt('export_data.json',outlier_json = "outliers_npmD.json",imagedir = "C:/Users/tituslim/Desktop/Allplots/")

