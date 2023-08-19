from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor

from datetime import datetime
import glob

#Functions to be used    
def edit_title(slide,title):
    slide.shapes.title.text = str(title)
    title = slide.shapes.title.text_frame.paragraphs[0]
    title.font.name = "Arial"
    title.font.size = Pt(54)
    title.font.bold = True

def addtxtbox(slide,header,comments,color=False):
    txBox = slide.shapes.add_textbox(left=Inches(8.68),top=Inches(2),height=Inches(0.8),width=Inches(2))
    tf = txBox.text_frame
    tf.text = header
    tf.paragraphs[0].font.name = "Arial"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True

    if comments == '':
        p.text = "Null"
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.bold = False
    else:
        #manipulate comments (outlier string)
        comment_lst = comments.split(', ')
        for comment in comment_lst: 
            p = tf.add_paragraph()
            p.text = comment
            p.font.name = "Arial"
            p.font.size = Pt(14)
            p.font.bold = False

    if color == True:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = RGBColor(100,200,100)

#def customtxtbox(slide): 
    #for shape in slide.shapes:
        #if shape.has_text_frame == True:


def plot_slide(prs,header=None,comments=None,title=None,image=None):
    layout = prs.slide_layouts[13]
    slide = prs.slides.add_slide(layout)
    edit_title(slide,title)
    slide.shapes.add_picture(image,height=Inches(5),width=Inches(8.34),left=Inches(0),top=Inches(2.2))
    addtxtbox(slide, header, comments)

    #Header format: '__ & __:' ; __ is variable(s) name
    #Comments format: 'Some Deviation/Comparable'

def transitn_slide(prs, title):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    edit_title(slide, title)
