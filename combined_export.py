import pandas as pd
# import numpy as np  #keep first, might be useful
import seaborn as sns
import matplotlib.pyplot as plt
from autoscale_visualization import scale, trendline, outliers_json
from sieve_outlier_anova import get_anova_score, stats, get_outlier, remove_outlier
from ksm_plot_v1 import kernel_smooth
# from export_genrl import exptppt #exclude first
import logging
import datetime as dt
import click
import os
import json
from pptx import Presentation
from pptx.dml.color import RGBColor
import glob
from export_funcs import edit_title, plot_slide, transitn_slide, addtxtbox
from jsonexport import exptppt
#from json_create import create_json

@click.command()
@click.argument('csv_file')
# @click.option('--operation', required=True, type=click.Choice(['autoscale', 'porconv'], case_sensitive=False), help='Choose to autoscale visualization or porconv comparison')
# , type=click.Choice(['1250-51 SLIT RECESS WET ETCH WEIGHT::MeasurementData::Product::MASS::DELTA_MASS - RAW_WAFER (Mean)', 'FINAL FUNCT PROD::WaferData::npnC', 'FINAL FUNCT PROD::WaferData::npmD', 'FINAL FUNCT PROD::WaferData::npuZQ']), help='Choose your variable of interest')
@click.option('--interest', required=True)
# interest leave as it is for a while, we may plan to just plot all interested variables wihtin the csv
@click.option('--ppt/--no-ppt', default=False, help='Do you want to export ppt?')
def main(csv_file, ppt, interest):  # removed operation option for now as it's not used

    # interest_probes = ['FINAL FUNCT PROD::WaferData::npmD',
    # 'FINAL FUNCT PROD::WaferData::npnC', 'FINAL FUNCT PROD::WaferData::npuZQ']
   
    if ppt == True:
        print('PPT')
        # exptppt('ppt_details.json', outlier_json="outliers_npmD.json", imagedir="C:\\Users\\tituslim\\Desktop\\Allplots\\")
    else:
        print('Nothing')


    start = dt.datetime.now()

    logger = logging.getLogger(__name__)
    logger.setLevel(logging.INFO)
    formatter = logging.Formatter('%(name)s - %(levelname)s - %(message)s')
    file_handler = logging.FileHandler('testing.log')
    file_handler.setFormatter(formatter)
    logger.addHandler(file_handler)
    logger.info('-' * 30)
    logger.info(f'run start: {start}')

    # will update later
    # datetime = "5450-51 SLIT RECESS WET ETCH::RunWaferData::ProcessEndDateTime"
    datetime = "5450-51 SLIT RECESS WET ETCH::RunData::ProcessEndDateTime"

    for i in range(9):
        # load data and remove NaN values, then sort them
        logger.info(f'reading csv file... {csv_file}')
        data = pd.read_csv(csv_file)
        data[datetime] = pd.to_datetime(data[datetime])
        logger.info(f'removing NA rows...')
        data = data.loc[data[datetime].notna()]
        data = data.loc[data[interest].notna()]

        # get por and conv
        por = data.loc[data["porconv"] == "por"]
        conv = data.loc[data["porconv"] == "conv"]
        por = por.sort_values(by=[datetime])
        conv = conv.sort_values(by=[datetime])

        # get stats of raw data, before removing outliers
        save_path = 'C:\\Users\\tituslim\\exportcode\\'
        file_name = 'stats_npmD.json'
        completeName = os.path.join(save_path, file_name)

        with open(completeName, 'w') as f:
            json.dump({"job_id": "1234"}, f)
            click.echo("stats.json created")

        click.echo('POR stats before outlier removal: [Mean, Std, Median]')
        stats(por, interest, "before", "por")
        click.echo('CONV stats before outlier removal: [Mean, Std, Median]')
        stats(conv, interest, "before", "conv")

        # kernel smooth
        logger.info(f'kernel smoothing in progress...')
        y_kvp = kernel_smooth(por[datetime], por[interest], x_datetime=True)
        y_kvc = kernel_smooth(conv[datetime], conv[interest], x_datetime=True)

        end = dt.datetime.now()
        runtime = (end - start).seconds
        logger.info(f'prep data and kernel smooth: {runtime}')

        # append the smoothed value to a new column for both por and conv
        por['y_kvp'] = y_kvp
        conv['y_kvc'] = y_kvc

        # sieve outliers
        logger.info("outlier sieving in progress...")
        data = get_anova_score(data, interest)
        cleaned_data = remove_outlier(data)
        cleaned_por = cleaned_data.loc[cleaned_data["porconv"] == "por"]
        cleaned_conv = cleaned_data.loc[cleaned_data["porconv"] == "conv"]
        cleaned_por = cleaned_por.sort_values(by=[datetime])
        cleaned_conv = cleaned_conv.sort_values(by=[datetime])

        # pass a json data of wafer id and lot id later from this csv, do later
        outliers = get_outlier(data)
        outliers.to_csv("outliers_data.csv")
        outliers_json(outliers)

        # stats after removing outliers
        click.echo('POR stats after outlier removal: [Mean, Std, Median]')
        stats(cleaned_por, interest, "after", "por")
        click.echo('CONV stats after outlier removal: [Mean, Std, Median]')
        stats(cleaned_conv, interest, "after", "conv")

        end = dt.datetime.now()
        runtime = (end - start).seconds
        logger.info(f'outliers sieving: {runtime}')

        #Read export_data.json and outlier_json
        #Json Data Parsing
        logger.info(f'JSON Data Parsing in progress...')
        f_1 = open("export_data.json",'r')
        export_data = json.loads(f_1.read())

        f_2 = open("outliers_npmD.json",'r')
        outlier_data = json.loads(f_2.read())

        #Create outlier_lst for input into ppt
        outlier_lst = [] #List of strings, string format = LotID: WaferId1, WaferId2,...
        for lotid,waferids in outlier_data.items():
            if len(waferids) > 1:
                one_line = lotid + ': '
                for waferid in waferids:
                    if waferid == waferids[-1]:
                        one_line += waferid 
                        outlier_lst.append(one_line)
                    else:
                        one_line += waferid + ','
            else:
                outlier_lst.append('{}: {}'.format(lotid,waferids[0]))
        
        outliers = '' #just 1 string = LotId: WaferId1, WaferId2,...
        for id in outlier_lst:
            if id == outlier_lst[-1]:
                outliers += id
            else:
                outliers += id + ', '

        end = dt.datetime.now()
        runtime = (end - start).seconds
        logger.info(f'JSON Data Parsing: {runtime}')

        #Making initial Slides
        #Create Presentation
        logger.info(f'Presentation Creation in progress...')
        prs = Presentation(export_data[0]["home"] + export_data[0]["template"])

        #Main Page, hardcoded(rarely changed)
        main_slide = prs.slides[0]
        date = dt.datetime.now()
        edit_title(main_slide,export_data[0]["title"])
        subtitle = main_slide.shapes.placeholders[25]
        subtitle.text = date.strftime("%m/%d/%Y, %H:%M:%S")

        #Main page color defaults (hardcoded change when needed)
        for shape in main_slide.shapes:
            if shape.shape_id == 2:
                shape.fill.background()
            if shape.shape_id == 9: 
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(150,50,100)
            if shape.shape_id == 10: 
                shape.fill.background()
            if shape.shape_id == 11: 
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(100,100,100)

        transitn_slide(prs, export_data[1]["title"])

        end = dt.datetime.now()
        runtime = (end - start).seconds
        logger.info(f'Presentation Creation: {runtime}')

        # kernel smooth for cleaned data
        logger.info(f'kernel smoothing for cleaned data in progress...')
        y_kvp = kernel_smooth(
            cleaned_por[datetime], cleaned_por[interest], x_datetime=True)
        y_kvc = kernel_smooth(
            cleaned_conv[datetime], cleaned_conv[interest], x_datetime=True)

        end = dt.datetime.now()
        runtime = (end - start).seconds
        logger.info(f'kernel smoothing for cleaned data: {runtime}')

        # append the smoothed value to a new column for both cleaned por and conv
        cleaned_por['y_kvp'] = y_kvp
        cleaned_conv['y_kvc'] = y_kvc

        ### Autoscale Visualization ###
        logger.info(f'autoscale visualization for cleaned data in progress...')
        mean = cleaned_data[datetime].mean()
        std = cleaned_data[datetime].std()
        new = cleaned_data.loc[abs(cleaned_data[datetime] - mean) < 2*std]

        min_date = new[datetime].min()
        max_date = new[datetime].max()

        cleaned_por = cleaned_por.loc[(cleaned_por[datetime] >= min_date) &
                                    (cleaned_por[datetime] <= max_date)]
        cleaned_conv = cleaned_conv.loc[(cleaned_conv[datetime] >= min_date) &
                                        (cleaned_conv[datetime] <= max_date)]

        #new = scale(new, interest)
        plt.figure(figsize=(10, 10))
        sns.set_theme()
        scatter = sns.scatterplot(data=new, x=datetime, y=interest,
                                hue="porconv", palette=dict(por="orange", conv="blue"))
        por_trend = trendline(cleaned_por[datetime], cleaned_por["y_kvp"], "por")
        conv_trend = trendline(
            cleaned_conv[datetime], cleaned_conv["y_kvc"], "conv")
        plt.savefig(export_data[0]["image_dir"] + "autoscaled_Visualization_no-outliers.png")
        files = glob.glob(export_data[0]["image_dir"] + "*.png")

        #Slide Creation + Figure Insertion
        x = 2
        png_dir = [dir for dir in files if export_data[x]['title'] in dir]
        plot_slide(prs,header = export_data[x]["header"],comments=outliers,title=export_data[x]["title"],image=png_dir[0])
        x += 1

        end = dt.datetime.now()
        runtime = (end - start).seconds
        logger.info(f'autoscale visualization with no outliers: {runtime}')

        # zoom in y axis
    
        new_loop = scale(new, interest)
        plt.figure(figsize=(10, 10))
        sns.set_theme()
        scatter = sns.scatterplot(data=new_loop, x=datetime, y=interest,
                              hue="porconv", palette=dict(por="orange", conv="blue"))
        por_trend = trendline(cleaned_por[datetime], cleaned_por["y_kvp"], "por")
        conv_trend = trendline(
        cleaned_conv[datetime], cleaned_conv["y_kvc"], "conv")
        name = "autoscaled_Visualization_no-outliers_zoomed-y" + str(i+1) + ".png"       
        plt.savefig(export_data[0]["image_dir"] + name)
        files = glob.glob(export_data[0]["image_dir"] + "*.png")
        
        #Slide Creation + Figure Insertion
        png_dir = [dir for dir in files if export_data[x]['title'] in dir]
        plot_slide(prs,header = export_data[x]["header"],comments=outliers,title=export_data[x]["title"],image=png_dir[0])
        x += 1

    end = dt.datetime.now()
    runtime = (end - start).seconds
    logger.info(
        f'autoscale visualization with no outliers zoomed y: {runtime}')

    #Presentation Closing slide
    logger.info(f'Presentation Closing in progress...')
    layout = prs.slide_layouts[17]
    prs.slides.add_slide(layout)
    prs.save("combined_10pngs.pptx")

    end = dt.datetime.now()
    runtime = (end - start).seconds
    logger.info(
        f'Presentation Closing: {runtime}')
    
    #Creation of ppt_details.json
    # logger.info(f'ppt_details.json creating...')
    # create_json(csv_file,interest_probes)

    end = dt.datetime.now()
    runtime = (end - start).seconds
    logger.info(f'run end: {end}')
    logger.info(f'total runtime: {runtime}')


if __name__ == '__main__':
    main()