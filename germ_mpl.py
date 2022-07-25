from pptx import Presentation
from export_funcs import edit_title
import json
import datetime as dt
import time
import logging
from cell_data import parse_and_export
import click


@click.command()
# @click.argument('csv_file')
# @click.option('--operation', required=True, type=click.Choice(['autoscale', 'porconv'], case_sensitive=False), help='Choose to autoscale visualization or porconv comparison')
# , type=click.Choice(['1250-51 SLIT RECESS WET ETCH WEIGHT::MeasurementData::Product::MASS::DELTA_MASS - RAW_WAFER (Mean)', 'FINAL FUNCT PROD::WaferData::npnC', 'FINAL FUNCT PROD::WaferData::npmD', 'FINAL FUNCT PROD::WaferData::npuZQ']), help='Choose your variable of interest_probes[0]')
# @click.option('--interest', required=True)
# interest leave as it is for a while, we may plan to just plot all
# interest variables wihtin the csv
# @click.option('--ppt/--no-ppt', default=False,
#               help='Do you want to export ppt?')
@profile
def main():
    start = dt.datetime.now()

    logger = logging.getLogger(__name__)
    logger.setLevel(logging.INFO)
    formatter = logging.Formatter('%(name)s - %(levelname)s - %(message)s')
    file_handler = logging.FileHandler('testing.log')
    file_handler.setFormatter(formatter)
    logger.addHandler(file_handler)
    logger.info('-' * 30)
    logger.info(f'run start: {start}')

    f = open('germ_compare.json', 'r')
    data = json.loads(f.read())

    prs = Presentation('empty.pptx')

    # Main Page, hardcoded(rarely changed)
    main_slide = prs.slides[0]
    date = dt.datetime.now()
    edit_title(main_slide, 'GerM Table Visualization')
    subtitle = main_slide.shapes.placeholders[25]
    subtitle.text = date.strftime("%m/%d/%Y, %H:%M:%S")

    i = 0
    # This variable keeps count of the number of rows (parameters) plotted
    rows_plotted = 0
    # This variable serves for the if loop condition later on
    rows_per_slide = 30
    # variable to count the # of loops (hence # of slides for a set of 4 tools)
    loop_counter = 0
    # calculate variables to generalise for # of rows in germ json
    quotient = len(data) // rows_per_slide
    logger.info(f'Data Parsing + Exporting in progress...')
    # Need to generalise for any number of parameters
    while rows_plotted <= (len(data) + 2):
        # plot the first (rows_plotted * quotient) number of rows
        rows_plotted += 30
        parse_and_export(prs, data, i, rows_plotted)
        i += 30
        loop_counter += 1
        if rows_plotted == rows_per_slide * quotient:
            if loop_counter == (quotient):
                print('loop is running')
                # assign i as old value of upper
                i = rows_per_slide * quotient
                print(i)
                # assign new value to upper
                rows_plotted = len(data)
                print(rows_plotted)
                parse_and_export(prs, data, i, rows_plotted)
                rows_plotted = len(data) + 3
                break
            break
        # Iterative variables updating;

        # plot the remaining number of rows

    end = dt.datetime.now()
    runtime = (end - start).seconds
    logger.info(f'Data Parsing + Exporting: {runtime}')

    # Last slide
    layout = prs.slide_layouts[17]
    slide = prs.slides.add_slide(layout)

    prs.save('sixrows_coloured.pptx')
    end = time.time()
    print("Runtime: {} seconds".format(runtime))


if __name__ == '__main__':
    main()
