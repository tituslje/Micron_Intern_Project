from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from export_funcs import edit_title
import pandas as pd
from datetime import datetime

master = pd.read_csv('dummy2.csv')
prs = Presentation('empty.pptx')

#Main Page, hardcoded(rarely changed)
main_slide = prs.slides[0]
date = datetime.now()
edit_title(main_slide,'GerM Table Visualization')
subtitle = main_slide.shapes.placeholders[25]
subtitle.text = date.strftime("%m/%d/%Y, %H:%M:%S")

layout = prs.slide_layouts[13]
slide = prs.slides.add_slide(layout)
edit_title(slide,'GerM Table')
r = len(master.index)
c = len(master.columns)
i=0
table1 = slide.shapes.add_table(r+1,c,Inches(0),Inches(1.25),Inches(14),Inches(3)).table
table1.columns.width = Inches(8)
for header in list(master.columns.values):
    while i <= (c-1):
        table1.cell(0,i).text = header
        i += 1
        break

x=1
y=0

for idx,sers in master.iterrows():
    for item in sers:
        table1.cell(x,y).text_frame.text = str(item)
        y += 1
        if x == r+1:
            break
        if y == c:
            y = 0
            x += 1
            break

for cell in table1.iter_cells():
    texts = cell.text_frame
    p = texts.paragraphs[0]
    run = p.add_run()
    font = run.font
    font.size = Pt(5)

prs.save('table_1.pptx')