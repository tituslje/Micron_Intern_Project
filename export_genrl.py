from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor

from datetime import datetime
import json
import os
from export_funcs import edit_title, addtxtbox, plot_slide, transitn_slide

#inputs accepted:
#1. template ppt file (needed as layouts are indexed from slide master)
#2. slide titles (transition, plots) (pass as list in order)
## list order: 
#[main,transition,probe plots...,transition,FD plots..,close]
#3. Image files list
##list order: [probe1,probe2,...,FD1,FD2..]
#3. # of probe plots, # of FD plots
#4. headers and comments for plot slides
#5. final ppt file save

def exptppt(data):
    #Create Presentation
    prs = Presentation(data[0]["template"])

    #Main Page, hardcoded(rarely changed)
    main_slide = prs.slides[0]
    date = datetime.now()
    edit_title(main_slide,data[0]["title"])
    subtitle = main_slide.shapes.placeholders[25]
    subtitle.text = date.strftime("%m/%d/%Y, %H:%M:%S")

    #Main page color defaults (hardcoded change when needed)
    for shape in main_slide.shapes:
        if shape.shape_id == 2:
            shape.fill.background()
        if shape.shape_id == 9: 
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0,100,100)
        if shape.shape_id == 10: 
            shape.fill.background()
        if shape.shape_id == 11: 
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(200,100,0)

    x = 2
    #First transition and Probe plots
    transitn_slide(prs, data[1]["title"])
    for i in range(data[0]["probe"]):
        plot_slide(prs,data[x]["header"],data[x]["comment"],data[x]["title"],data[x]["image"])
        x += 1
    
    #Second transition, FD plots, and last slide
    transitn_slide(prs,data[x]["title"])
    for i in range(data[0]["FD"]):
        if data[0]["FD"] == 0: #No FD Plots
            quit()
        if data[0]["FD"] == 1: #Only 1 FD Plot
            i = 1
            plot_slide(prs,data[x+i]["header"],data[x+i]["comment"],data[x+i]["title"],data[x+i]["image"])
            quit()
        if data[0]["FD"] > 1: #More than 1 FD Plot (like Probe)
            i += 1 #Increment i by 1 to avoid recreating past slide
            plot_slide(prs,data[x+i]["header"],data[x+i]["comment"],data[x+i]["title"],data[x+i]["image"])

    layout = prs.slide_layouts[17]
    prs.slides.add_slide(layout)
    
    prs.save("json.pptx")

f = open('export_data.json','r')
data = json.loads(f.read())
exptppt(data)