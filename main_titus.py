import click
import datetime as dt
import glob
import json
import logging
import matplotlib.pyplot as plt
import os
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
import seaborn as sns
from autoscale_visualization import scale, trendline
from export_funcs import edit_title, plot_slide, transitn_slide
from json_create import create_json
from ksm_plot_v1 import kernel_smooth
from sieve_outlier_anova import get_anova_score, stats, remove_outlier, get_outlier, outliers_to_json


@click.command()
@click.argument('csv_file')
@click.option('--ppt/--no-ppt', default=False,
              help='Do you want to export ppt?')
@profile
def main(csv_file, ppt):  

    interest_probes = [
        'FINAL FUNCT PROD::WaferData::npuZQ',
        'FINAL FUNCT PROD::WaferData::npnC',
        'FINAL FUNCT PROD::WaferData::npmD']

    start = dt.datetime.now()

    logger = logging.getLogger(__name__)
    logger.setLevel(logging.INFO)
    formatter = logging.Formatter('%(name)s - %(levelname)s - %(message)s')
    file_handler = logging.FileHandler('testing.log')
    file_handler.setFormatter(formatter)
    logger.addHandler(file_handler)
    logger.info('-' * 30)
    logger.info(f'run start: {start}')

    if ppt:
        print('PPT')
        create_json(csv_file, interest_probes)
        f_1 = open("ppt_details.json", 'r')
        export_data = json.loads(f_1.read())

        # Making initial Slides
        # Create Presentation
        logger.info(f'Presentation Creation in progress...')
        prs = Presentation(export_data[0]["home"] + export_data[0]["template"])

        # Main Page, hardcoded(rarely changed)
        main_slide = prs.slides[0]
        date = dt.datetime.now()
        edit_title(main_slide, export_data[0]["title"])
        subtitle = main_slide.shapes.placeholders[25]
        subtitle.text = date.strftime("%m/%d/%Y, %H:%M:%S")

        # Main page color defaults (hardcoded change when needed)
        for shape in main_slide.shapes:
            if shape.shape_id == 2:
                shape.fill.background()
            if shape.shape_id == 9:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(150, 50, 100)
            if shape.shape_id == 10:
                shape.fill.background()
            if shape.shape_id == 11:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(100, 100, 100)

        transitn_slide(prs, export_data[1]["title"])

        end = dt.datetime.now()
        runtime = (end - start).seconds
        logger.info(f'Presentation Creation: {runtime}')
        # exptppt('ppt_details.json', outlier_json="outliers_npmD.json", imagedir="C:\\Users\\tituslim\\Desktop\\Allplots\\")
    else:
        print('Nothing')

    datetime = "5450-51 SLIT RECESS WET ETCH::RunData::ProcessEndDateTime"

    for interest in interest_probes:
        # load data and remove NaN values, then sort them
        logger.info(f'reading csv file... {csv_file}')
        data = pd.read_csv(csv_file)
        data[datetime] = pd.to_datetime(data[datetime])
        logger.info(f'removing NA rows...')
        data = data.loc[data[datetime].notna()]
        data = data.loc[data[interest].notna()]

        # get por and conv
        por = data.loc[data["porconv"] == "por"]
        conv = data.loc[data["porconv"] == "conv"]
        por = por.sort_values(by=[datetime])
        conv = conv.sort_values(by=[datetime])

        # get stats of raw data, before removing outliers
        save_path = 'C:\\Users\\tituslim\\exportcode\\'
        file_name = 'stats.json'
        completeName = os.path.join(save_path, file_name)

        with open(completeName, 'w') as f:
            json.dump({"job_id": "1234"}, f)

            click.echo("stats.json created")

        click.echo('POR stats before outlier removal: [Mean, Std, Median]')
        stats(por, interest, "before", "por")
        click.echo('CONV stats before outlier removal: [Mean, Std, Median]')
        stats(conv, interest, "before", "conv")

        # kernel smooth
        logger.info(f'kernel smoothing in progress...')
        y_kvp = kernel_smooth(por[datetime], por[interest], x_datetime=True)
        y_kvc = kernel_smooth(conv[datetime], conv[interest], x_datetime=True)

        end = dt.datetime.now()
        runtime = (end - start).seconds
        logger.info(f'prep data and kernel smooth: {runtime}')

        # append the smoothed value to a new column for both por and conv
        por['y_kvp'] = y_kvp
        conv['y_kvc'] = y_kvc

        # Read outlier_json
        # Json Data Parsing
        logger.info("outlier data parsing...")
        file_name1 = 'outliers.json'
        outliers_json_1 = os.path.join(save_path, file_name1)
        with open(outliers_json_1, 'w') as f:
            dump = dict(zip(interest_probes, [None] * len(interest_probes)))
            json.dump(dump, f)
            click.echo('outliers.json created')

        # sieve outliers
        logger.info("outlier sieving in progress...")
        data = get_anova_score(data, interest)
        cleaned_data = remove_outlier(data)
        cleaned_por = cleaned_data.loc[data["porconv"] == "por"]
        cleaned_conv = cleaned_data.loc[data["porconv"] == "conv"]
        # cleaned_por = cleaned_por.sort_values(by=[datetime])
        # cleaned_conv = cleaned_conv.sort_values(by=[datetime])

        # pass a json data of wafer id and lot id later from this csv, do later
        outliers = get_outlier(data)
        outliers.to_csv("outliers_data.csv")
        outliers_to_json(outliers, interest)

        # stats after removing outliers
        click.echo('POR stats after outlier removal: [Mean, Std, Median]')
        stats(por, interest, "after", "por")
        click.echo('CONV stats after outlier removal: [Mean, Std, Median]')
        stats(conv, interest, "after", "conv")

        end = dt.datetime.now()
        runtime = (end - start).seconds
        logger.info(f'outliers sieving: {runtime}')

        f_2 = open("outliers.json", 'r')
        outlier_data = json.loads(f_2.read())

        # Create outlier_lst for input into ppt
        # List of strings, string format = LotID: WaferId1, WaferId2,...
        outlier_lst = []
        outliers_probe = ''
        for probe, values in outlier_data.items():
            if interest == probe:
                for lotid, waferids in values.items():
                    if len(waferids) > 1:
                        one_line = lotid + ': '
                        for waferid in waferids:
                            if waferid == waferids[-1]:
                                one_line += waferid
                                outlier_lst.append(one_line)
                            else:
                                one_line += waferid + ','
                    else:
                        outlier_lst.append('{}: {}'.format(lotid, waferids[0]))

        # just 1 string = LotId: WaferId1, WaferId2,...
        for id in outlier_lst:
            if id == outlier_lst[-1]:
                outliers_probe += id
            else:
                outliers_probe += id + ', '

        end = dt.datetime.now()
        runtime = (end - start).seconds
        logger.info(f'Outlier Data Parsing: {runtime}')

        # kernel smooth for cleaned data
        # logger.info(f'kernel smoothing for cleaned data in progress...')
        # y_kvp = kernel_smooth(
        #     por[datetime], por[interest], x_datetime=True)
        # y_kvc = kernel_smooth(
        #     conv[datetime], conv[interest], x_datetime=True)

        # end = dt.datetime.now()
        # runtime = (end - start).seconds
        # logger.info(f'kernel smoothing for cleaned data: {runtime}')

        # append the smoothed value to a new column for both cleaned por and conv
        # por['y_kvp'] = y_kvp
        # conv['y_kvc'] = y_kvc

        ### Autoscale Visualization ###
        logger.info(f'autoscale visualization for cleaned data in progress...')
        mean = data[datetime].mean()
        std = data[datetime].std()
        new = data.loc[abs(data[datetime] - mean) < 2 * std]

        min_date = new[datetime].min()
        max_date = new[datetime].max()

        por = por.loc[(por[datetime] >= min_date) &
                      (por[datetime] <= max_date)]
        conv = conv.loc[(conv[datetime] >= min_date) &
                        (conv[datetime] <= max_date)]

        #new = scale(new, interest)
        plt.figure(figsize=(10, 10))
        sns.set_theme()
        scatter = sns.scatterplot(
            data=new,
            x=datetime,
            y=interest,
            hue="porconv",
            palette=dict(
                por="orange",
                conv="blue"))
        por_trend = trendline(por[datetime], por["y_kvp"], "por")
        conv_trend = trendline(
            conv[datetime], conv["y_kvc"], "conv")
        name_lst = interest.split('::')
        # #Create figure (png file) directory
        plot_name = name_lst[-1] + "_scaled_no-outliers"
        with open('ppt_details.json', 'w') as f:
            # load = json.load(f)
            title = plot_name
            # load.append({"title":title,"header": "Por & Conv Outliers:","image_dir":"C:\\Users\\tituslim\\Desktop\\Allplots\\"})
            json.dump({"title": title,
                       "header": "Conv Outliers:",
                       "image_dir": "C:\\Users\\tituslim\\Desktop\\Allplots\\"},
                      f,
                      indent=4)

        f_1 = open("ppt_details.json", 'r')
        export_data = json.loads(f_1.read())
        # #Save figure using updated json
        plt.savefig(export_data["image_dir"] + plot_name + ".png")
        # #update files directory to include new png saved
        files = glob.glob(export_data["image_dir"] + "*.png")

        # Slide Creation + Figure Insertion
        png_dir = [dir for dir in files if export_data['title'] in dir]
        plot_slide(
            prs,
            header=export_data["header"],
            comments=outliers_probe,
            title=export_data["title"],
            image=png_dir[0])
        os.remove(export_data["image_dir"] + plot_name + ".png")

        end = dt.datetime.now()
        runtime = (end - start).seconds
        logger.info(f'autoscale visualization with no outliers: {runtime}')

        # zoom in y axis
        new_loop = scale(new, interest)
        plt.figure(figsize=(10, 10))
        sns.set_theme()
        scatter = sns.scatterplot(
            data=new_loop,
            x=datetime,
            y=interest,
            hue="porconv",
            palette=dict(
                por="orange",
                conv="blue"))
        por_trend = trendline(por[datetime], por["y_kvp"], "por")
        conv_trend = trendline(
            conv[datetime], conv["y_kvc"], "conv")
        # #Create figure (png file) directory
        plot_name = name_lst[-1] + "_scaled_no-outliers_zoomed-y"
        with open('ppt_details.json', 'w') as f:
            # load = json.load(f)
            title = plot_name
            json.dump({"title": title,
                       "header": "Conv Outliers:",
                       "image_dir": "C:\\Users\\tituslim\\Desktop\\Allplots\\"},
                      f,
                      indent=4)

        f_1 = open("ppt_details.json", 'r')
        export_data = json.loads(f_1.read())
        # #Save figure using updated json
        plt.savefig(export_data["image_dir"] + plot_name + ".png")
        # #update files directory to include new png saved
        files = glob.glob(export_data["image_dir"] + "*.png")

        # Slide Creation + Figure Insertion
        png_dir = [dir for dir in files if export_data['title'] in dir]
        plot_slide(
            prs,
            header=export_data["header"],
            comments=outliers_probe,
            title=export_data["title"],
            image=png_dir[0])
        os.remove(export_data["image_dir"] + plot_name + ".png")

    end = dt.datetime.now()
    runtime = (end - start).seconds
    logger.info(
        f'autoscale visualization with no outliers zoomed y: {runtime}')

    # Presentation Closing slide
    logger.info(f'Presentation Closing in progress...')
    layout = prs.slide_layouts[17]
    prs.slides.add_slide(layout)
    prs.save("new.pptx")

    end = dt.datetime.now()
    runtime = (end - start).seconds
    logger.info(
        f'Presentation Closing: {runtime}')

    end = dt.datetime.now()
    runtime = (end - start).seconds
    logger.info(f'run end: {end}')
    logger.info(f'total runtime: {runtime}')


if __name__ == '__main__':
    main()
