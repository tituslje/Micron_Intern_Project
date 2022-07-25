from pptx.util import Inches
from export_funcs import edit_title
import os
import matplotlib.pyplot as plt
from matplotlib.colors import is_color_like
import pandas as pd

# takes presentation object, json data object,


def parse_and_export(prs, json_data, i, upper):
    cols = []
    rows = []
    cell_txt = []
    cell_txt_no_colour = []

    # Accumulate cols to have a list of column labels
    while len(cols) < 8:
        for k, v in json_data[0].items():
            if '(' in k and k not in ["ParamName", "ParamModule", "ParamSeq"]:
                if '#' in v and v not in cols:  # append column labels without colour code
                    column = v[0:len(v)-7]  # colour code strings have length 7
                    cols.append(column)
                if '#' not in v and v not in cols:  # colour code absent;no removal needed
                    cols.append(v)

    # Label the columns as reference (ref) or Point of Interest (POI)
    col_annotated = []
    for col in cols:
        if 'A1C' in col:
            col_annotated.append(str(col) + str(' (REF)'))
        if 'A72' in col:
            col_annotated.append(str(col) + str(' (POI)'))

    # Loop that accumulates data, and exports them for 4 rows each time
    while i < upper:
        # Initialise temporary arrays containing row data (left to right) which resets every loop
        temp_no_colour = []
        temp = []

        # Accumulate row labels, 1 row label per loop
        rows.append(json_data[i]["ParamName"])

        # Accumulate cell data in list, colour codes remain
        for k, v in json_data[i].items():
            if 'Param' not in k:  # Filter out tool names as key variables
                if v == "":  # Account for empty dictionary value in key
                    temp.append('NAN')
                if v != "":  # non empty dictionary value in key
                    temp.append(v)

        # Accumulate cell data in list, and remove colour codes from individual cells
        for k, v in json_data[i].items():
            if 'Param' not in k:  # filter out tool names as key variables
                if v == "":  # empty dictionary value in key
                    temp_no_colour.append('NAN')
                if v != "":  # non empty dictionary value in key
                    # long string with colour code
                    if '#' in v:
                        if len(v) > 10:
                            value = v[0:len(v)-7]
                            length = len(value)
                            if len(value) > 12:
                                # new_val = value[0:length//2] + \
                                #     '\n' + value[length//2:len(v)]
                                temp_no_colour.append(value)
                            else:
                                temp_no_colour.append(value)
                        if len(v) <= 10:
                            # short string is not entirely a colour code
                            if is_color_like(v) == False:
                                temp_no_colour.append(v[0:len(v)-7])
                            # short string is colour code
                            if is_color_like(v) == True:
                                temp_no_colour.append('')
                    # long string no colour code
                    if '#' not in v and len(v) > 10:
                        # length = len(v)
                        # new_val = v[0:length//2] + \
                        #     '\n' + v[length//2:len(v)]
                        temp_no_colour.append(v)
                    # short string no colour code
                    if '#' not in v and len(v) <= 10:
                        temp_no_colour.append(v)

        cell_txt.append(temp)
        cell_txt_no_colour.append(temp_no_colour)
        i += 1

    # Condition so while loop is not exited; 2 Dataframes made at last loop condition.
    if i == (upper):
        # Generate dataframe with colour codes in cells. This is used for extracting colour codes.
        j = 0
        master = pd.DataFrame(columns=col_annotated, index=rows)
        for row_data in cell_txt:
            master.loc[rows[j]] = row_data
            j += 1

        # Generate dataframe without colour codes. This is used for plotting the table.
        j = 0
        master_no_colour = pd.DataFrame(columns=col_annotated, index=rows)
        for row_data in cell_txt_no_colour:
            master_no_colour.loc[rows[j]] = row_data
            j += 1

        # Generating a list of column labels, sorted in the desired order:
        # reference column label followed by comparison column label
        cols_to_compare = [col for col in master_no_colour.columns]
        cols_to_compare.sort()

        # A while loop that splits dataframe to show 2 tools per slide at the end of each loop.
        first = 0
        last = 2
        while last <= 8:
            # Temporary dataframe which extracts the 2 column labels (or 2 tools) to be compared.
            # This changes every loop to cover all 8 labels for the same 4 rows (parameter names).
            # 4 slides --> all 8 column labels (or tools), 4 rows (or parameter names)
            temp_df = master[cols_to_compare[first:last]]

            # Extract colour codes from dataframe master in a list.
            # Accumulation is from left to right, top to bottom convention.
            colors_all = []
            for idx, row in temp_df.iterrows():
                colors = []
                for cell in row:
                    if len(colors) < 2:
                        if '#' in cell:
                            colors.append(cell[-7:])
                        if '#' not in cell:
                            colors.append('1')
                for color in colors:
                    if '\n' in color:
                        color.replace('\n', '#')
                colors_all.append(colors)

            # Temporary dataframe which extracts the 2 column labels to be compared.
            # This changes every loop to cover all 8 labels for the same 4 rows (parameter names).
            # 4 slides --> all 8 column labels (or tools), 4 rows (or parameter names)
            temp_df_no_colour = master_no_colour[cols_to_compare[first:last]]

            # Plotting of Table
            fig, ax = plt.subplots()
            # Final table is coloured and colour codes in cells are removed,
            # with the help of 2 dataframes.
            the_table = pd.plotting.table(
                ax, temp_df_no_colour, loc='center', cellLoc='center', cellColours=colors_all, colWidths=[1, 1])
            # the_table = plt.table(cellText=cell_txt,cellColours=None,rowLabels = rows,colLabels = cols,loc='center',cellLoc='center')

            # Autoscale table cells to fit to word length
            the_table.auto_set_font_size(True)
            the_table.set_fontsize(12)
            # the_table.auto_set_column_width(col=list(range(len(cols))))
            # the_table.auto_set_column_width(col=list(range(2)))
            the_table.scale(1, 4)

            # The table portion of graph is isolated.
            ax.get_xaxis().set_visible(False)
            ax.get_yaxis().set_visible(False)
            plt.box(on=None)

            plt.savefig("temp_store.png", bbox_inches='tight')
            fig.clear()
            plt.close('all')

            # Sub-slides editing
            layout = prs.slide_layouts[13]
            slide = prs.slides.add_slide(layout)
            edit_title(slide, 'GerM Table')

            # slide.shapes.add_picture only accepts image files, or path to an image (which still needs a saved image file)
            slide.shapes.add_picture("temp_store.png", height=Inches(
                5.36), width=Inches(12.4), left=Inches(0.5), top=Inches(1.45))

            # Remove image file from directory to save local memory space.
            os.remove("C:\\Users\\tituslim\\germtable\\temp_store.png")

            # Update iterative variables to go to next 2 column labels (tools) to be plotted.
            first += 2
            last += 2
